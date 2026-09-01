# -*- coding: utf-8 -*-
"""iCAN 产品说明 PPT 生成（13 页，16:9，古风配色）。

前置：先运行 record_demo.py 产出原型截图。
运行：python scripts/make_ppt.py → submission/千年晤旅_产品说明.pptx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from config import PROJECT_ROOT

SUB = PROJECT_ROOT / "submission"
SHOT_DIR = SUB / "screenshots"
CHAR_DIR = PROJECT_ROOT / "assets" / "characters"

PAPER = RGBColor(0xF7, 0xF1, 0xE3)
PAPER2 = RGBColor(0xEF, 0xE6, 0xD0)
INK = RGBColor(0x3A, 0x32, 0x26)
INK_L = RGBColor(0x8A, 0x7C, 0x66)
SEAL = RGBColor(0xC0, 0x39, 0x2B)
GOLD = RGBColor(0xB0, 0x8D, 0x4F)
WHITE = RGBColor(0xFF, 0xFD, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _set_font(run, size=16, color=INK, bold=False, name="微软雅黑"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = PAPER
    return s


def add_box(slide, x, y, w, h, lines, size=15, color=INK, bold=False,
            align=PP_ALIGN.LEFT, fill=None, line_color=None, title_size=None):
    """lines: str 或 [(text, {style overrides})]。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = [(lines, {})]
    for i, (text, style) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = style.get("align", align)
        run = p.add_run()
        run.text = text
        _set_font(run, style.get("size", title_size or size),
                  style.get("color", color), style.get("bold", bold))
    if fill is not None or line_color is not None:
        # 用形状代替背景
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y), Inches(w), Inches(h))
        if fill is not None:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        else:
            shp.fill.background()
        if line_color is not None:
            shp.line.color.rgb = line_color
            shp.line.width = Pt(1.2)
        else:
            shp.line.fill.background()
        shp._element.addnext(tb._element)  # 文本置于形状之上
    return tb


def title_bar(slide, text, sub=None):
    add_box(slide, 0.55, 0.32, 12.2, 0.75, text, size=30, bold=True)
    d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.05),
                               Inches(12.13), Pt(2.2))
    d.fill.solid()
    d.fill.fore_color.rgb = GOLD
    d.line.fill.background()
    if sub:
        add_box(slide, 0.6, 1.15, 12, 0.4, sub, size=13, color=INK_L)


def footer(slide, idx):
    add_box(slide, 12.2, 7.05, 0.9, 0.35, f"{idx:02d}", size=12, color=INK_L,
            align=PP_ALIGN.RIGHT)
    add_box(slide, 0.55, 7.05, 6, 0.35, "千年晤旅 · 沉浸式历史人文游历交互平台",
            size=10, color=INK_L)


def arrow_down(slide, x, y, w=0.35, h=0.3):
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = GOLD
    a.line.fill.background()


def flow_box(slide, x, y, w, h, text, fill=PAPER2, bold=False, size=13):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = INK_L
    shp.line.width = Pt(1)
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    _set_font(r, size=size, bold=bold)
    return shp


def shot_available(name):
    return (SHOT_DIR / f"{name}.png").exists()


# ==================== S1 封面 ====================
s = new_slide()
add_box(s, 2.5, 2.1, 8.3, 1.5, "千年晤旅", size=72, bold=True, align=PP_ALIGN.CENTER)
add_box(s, 2.5, 3.55, 8.3, 0.6, "沉浸式历史人文游历交互平台", size=30,
        color=INK_L, align=PP_ALIGN.CENTER)
seal = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.07), Inches(4.35),
                          Inches(1.2), Inches(1.2))
seal.fill.solid()
seal.fill.fore_color.rgb = SEAL
seal.line.color.rgb = WHITE
seal.line.width = Pt(2)
stf = seal.text_frame
sp = stf.paragraphs[0]
sp.alignment = PP_ALIGN.CENTER
sr = sp.add_run()
sr.text = "晤旅"
_set_font(sr, size=28, color=WHITE, bold=True)
add_box(s, 2.5, 5.85, 8.3, 0.5, "与古人同游故地 · 听亲历者讲往事", size=18,
        color=INK, align=PP_ALIGN.CENTER)
add_box(s, 2.5, 6.6, 8.3, 0.45, "iCAN 大学生创新创业大赛 · 智慧文创赛道", size=14,
        color=INK_L, align=PP_ALIGN.CENTER)

# ==================== S2 目录 ====================
s = new_slide()
title_bar(s, "目录")
items = ["01 背景与行业痛点", "02 国内外研究现状", "03 方案设计", "04 核心功能",
         "05 业务流程", "06 人物形象设计", "07 系统架构", "08 原型展示",
         "09 项目创新点", "10 不足与未来迭代", "11 总结致谢"]
for i, it in enumerate(items):
    col, row = i % 2, i // 2
    x = 0.9 + col * 6.0
    y = 1.75 + row * 0.82
    flow_box(s, x, y, 5.5, 0.62, it, size=16)

# ==================== S3 背景痛点 ====================
s = new_slide()
title_bar(s, "背景与行业痛点", "传统文旅产品“重信息、轻体验”，游客面对古迹时“知其然而不知其所以然”")
cards = [
    ("打卡同质化", "主流旅游平台提供景点列表与攻略，路线千人一面，旅行沦为“拍照打卡流水线”。"),
    ("路线没有故事", "路线只给出“去哪儿”，缺少历史情感与人物故事，游客不知道看什么历史。"),
    ("读不懂古迹", "古迹背后的历史心境难以共鸣——游客看不懂一块碑、一座祠承载的悲欢。"),
    ("科普生硬百科", "传统文化科普以百科词条式输出为主，缺乏代入感，难以吸引年轻人。"),
]
for i, (t, d) in enumerate(cards):
    x = 0.9 + (i % 2) * 6.0
    y = 1.9 + (i // 2) * 2.3
    flow_box(s, x, y, 5.5, 0.55, t, bold=True, size=16)
    add_box(s, x + 0.25, y + 0.65, 5.0, 1.5, d, size=13, color=INK_L)
add_box(s, 0.9, 6.6, 11.5, 0.5, "痛点核心：游客需要一个“懂历史、有情感、可对话”的同行者。",
        size=16, bold=True, color=SEAL)

# ==================== S4 国内外现状 ====================
s = new_slide()
title_bar(s, "国内外研究现状", "数字文旅兴起，但“人物第一人称 + 成套游历闭环”仍是空白")
rows = [
    ("类别", "代表形态", "局限"),
    ("传统旅游平台", "景点列表 / 攻略社区 / 智能推荐", "信息堆砌，缺历史情感与人物故事"),
    ("数字文博", "AR/VR 导览、数字博物馆", "以展品为中心，静态讲解，无个性化"),
    ("大模型陪伴应用", "通用 AI 聊天 / 虚拟角色", "非文旅垂直场景，无地点与路线闭环"),
    ("千年晤旅", "地点→人物→偏好融合路线→故地对话→打卡合影→档案 + 流式对话/语音输入", "第一人称叙事 + 完整游历链路一体化"),
]
for r, (c1, c2, c3) in enumerate(rows):
    y = 1.75 + r * 1.02
    fills = [PAPER2, WHITE, WHITE, WHITE, WHITE]
    bolds = [True, True, False, False, False]
    flow_box(s, 0.9, y, 2.6, 0.85, c1, fill=fills[r], bold=bolds[r], size=14)
    flow_box(s, 3.6, y, 4.4, 0.85, c2, fill=fills[r], bold=bolds[r], size=13)
    flow_box(s, 8.1, y, 4.4, 0.85, c3, fill=fills[r], bold=bolds[r], size=13)
    if r == 4:  # 千年晤旅行：描红强调
        s.shapes[-3].line.color.rgb = SEAL
        s.shapes[-2].line.color.rgb = SEAL
        s.shapes[-1].line.color.rgb = SEAL

# ==================== S5 方案设计 ====================
s = new_slide()
title_bar(s, "方案设计", "拒绝生硬百科——让历史人物亲口讲述自己的故事")
principles = [
    ("第一人称叙事", "全部历史内容以人物“我”的口吻输出，还原事件发生时的内心想法与情绪"),
    ("双向融合路线", "AI 同时读取【用户游玩喜好】与【人物真实生平爱好经历】，成套编排行程"),
    ("双模式路线", "模式A 人物视角优先 / 模式B 双向融合，自由切换"),
    ("双使用场景", "实地文旅出行全程陪游 + 日常休闲与古人聊天学习"),
]
for i, (t, d) in enumerate(principles):
    x = 0.9 + (i % 2) * 6.0
    y = 1.8 + (i // 2) * 2.35
    flow_box(s, x, y, 5.5, 0.55, t, fill=WHITE, bold=True, size=15)
    add_box(s, x + 0.25, y + 0.62, 5.0, 1.5, d, size=13, color=INK_L)

# ==================== S6 核心功能 ====================
s = new_slide()
title_bar(s, "核心功能", "九大功能模块 + 三大年轻化玩法")
funcs = [
    ("历史人物搜索/推荐", "按姓名、身份、事迹检索 17 位人物；课本热点、网络热门分组推荐"),
    ("地点检索 + GPS", "输入城市/景点匹配在此生活游历过的历史人物；支持古称别名、附近故地定位"),
    ("偏好采集", "古迹/美食/山水/节奏等多项选择 + 自由想法 + 天数 + 双模式路线"),
    ("AI 成套路线", "多天成套行程，每站绑定人物故事与开场白；AI 不可用时内置足迹拼装降级"),
    ("故地重游对话", "人物第一人称讲此地往事与当时心境；AI 在线时回复逐字流式浮现"),
    ("双地图点亮", "古风画卷 SVG 地图 + 实景交互地图（点击站点直达对话），完成对话解锁打卡点亮"),
    ("AI 同游合影 + 随笔", "上传自拍合成古风合影；写游历随笔（可请古人润色，亦流式输出）"),
    ("个人游历档案", "路线、点亮足迹、合影、随笔全部落盘，随时回看并支持导出长图/行程单/分享卡片（真二维码）"),
    ("单人对话 + 群聊", "第一人称闲谈（支持浏览器语音输入）；跨时代群聊并发生成、逐条浮现"),
]
for i, (t, d) in enumerate(funcs):
    x = 0.9 + (i % 3) * 4.0
    y = 1.7 + (i // 3) * 1.32
    flow_box(s, x, y, 3.75, 0.5, t, bold=True, size=13.5)
    add_box(s, x + 0.12, y + 0.55, 3.5, 0.68, d, size=11.5, color=INK_L)
add_box(s, 0.9, 5.85, 11.5, 0.45, "年轻化玩法：🎯 古今人格测试 · 🎁 古人赠藏头诗 · 🎫 时空票根收藏",
        size=15, bold=True, color=SEAL)

# ==================== S7 业务流程 ====================
s = new_slide()
title_bar(s, "业务流程", "旅游完整链路一体化串联，数据环环传递")
steps = ["① 输入旅游地点", "② 检索匹配当地历史人物", "③ 选定一位同行人物",
         "④ 填写游玩偏好/天数/模式", "⑤ AI 融合生成成套路线", "⑥ 逐站故地重游对话",
         "⑦ 解锁打卡·地图点亮", "⑧ AI同游合影·写随笔", "⑨ 收入个人游历档案"]
for i, st in enumerate(steps):
    y = 1.55 + i * 0.575
    flow_box(s, 0.7, y, 4.6, 0.45, st, fill=(WHITE if i % 2 else PAPER2), size=12.5)
    if i < len(steps) - 1:
        arrow_down(s, 2.9, y + 0.44, w=0.22, h=0.14)
add_box(s, 0.7, 6.75, 4.6, 0.4, "一条链路 · 九步闭环", size=13, bold=True, color=SEAL,
        align=PP_ALIGN.CENTER)
# 右侧：日常聊天分支 + 玩法
flow_box(s, 6.6, 1.55, 5.9, 0.55, "日常聊天流程（独立分支）", bold=True, size=14)
flow_box(s, 6.9, 2.35, 5.3, 0.45, "搜索/推荐历史人物", size=12.5)
arrow_down(s, 9.4, 2.79, w=0.22, h=0.14)
flow_box(s, 6.9, 2.95, 5.3, 0.45, "单人第一人称对话 / 创建跨时代群聊", size=12.5)
arrow_down(s, 9.4, 3.39, w=0.22, h=0.14)
flow_box(s, 6.9, 3.55, 5.3, 0.45, "与古人自由聊天 / 古人互辩思想碰撞", size=12.5)
add_box(s, 6.6, 4.35, 5.9, 0.5, "年轻化玩法", size=16, bold=True, color=SEAL)
for j, (t, d) in enumerate([("🎯 古今人格测试", "六题测出你的古代人格"),
                            ("🎁 古人赠藏头诗", "名字嵌入诗中，落款古人"),
                            ("🎫 时空票根", "旅程纪念票根收藏")]):
    y = 5.0 + j * 0.62
    flow_box(s, 6.9, y, 2.2, 0.48, t, bold=True, size=12)
    add_box(s, 9.25, y + 0.06, 3.3, 0.4, d, size=11.5, color=INK_L)

# ==================== S8 人物设计 ====================
s = new_slide()
title_bar(s, "人物形象设计", "史料真实可考 · 国风写实古风 · 多元身份")
demos = [("su_shi", "苏轼 · 北宋文学家"), ("li_bai", "李白 · 唐文学家"),
         ("li_qingzhao", "李清照 · 宋文学家"), ("bi_sheng", "毕昇 · 北宋发明家"),
         ("zhuge_liang", "诸葛亮 · 蜀汉政治家")]
for i, (pid, label) in enumerate(demos):
    x = 0.75 + i * 2.42
    p = CHAR_DIR / f"{pid}.png"
    if p.exists():
        s.shapes.add_picture(str(p), Inches(x), Inches(1.7), height=Inches(3.6))
    add_box(s, x - 0.2, 5.35, 2.8, 0.5, label, size=12, bold=True, align=PP_ALIGN.CENTER)
add_box(s, 0.75, 6.15, 11.8, 0.85,
        [("规范：仅选史书可靠记载的真实人物；美术贴合朝代服饰器物，拒绝Q版魔改；", {"size": 13}),
         ("人物库共 17 位（文学家/发明家/工匠/科学家/政治家），每人配置性格/口吻/事迹/爱好/名句 + 各地点第一人称故事文本（AI 素材与离线降级内容源）；", {"size": 13}),
         ("Demo 使用静态透明 PNG 半身图，动态数字人列入未来迭代。", {"size": 13})])

# ==================== S9 系统架构 ====================
s = new_slide()
title_bar(s, "系统架构", "轻量化 Python 原型 · 单入口状态机路由 · AI 带离线降级保障")
layers = [
    ("呈现层", "Streamlit 单入口 app.py + session_state 状态机路由（13 个页面视图）· 古风主题样式 + 移动端断点适配"),
    ("核心层", "路线编排 route_builder · 提示词模板 · 离线剧本 scripts · 双地图（SVG 画卷 + pydeck 实景交互）· GPS 定位 · 流式输出 · Web Speech 语音输入 · PIL 合影/票根 · SQLite 持久化"),
    ("数据层", "人物库 17 位（第一人称文本）· 地点库 37 处/24 城（真实经纬度 + 古称别名）· 双向索引 · 本地素材（立绘/背景）"),
    ("服务层", "大模型对话接口（OpenAI 兼容协议，DeepSeek）· 结构化 JSON 输出 · 流式生成 · 确定性缓存（同参数零重复付费）· 超时重试"),
    ("保障层", "AI_DISABLED 一键离线演示：内置路线/关键词应答/群聊剧本（5 组 + 通用台词），全链路闭环不受影响"),
]
for i, (t, d) in enumerate(layers):
    y = 1.7 + i * 1.02
    flow_box(s, 0.9, y, 1.7, 0.82, t, fill=PAPER2, bold=True, size=15)
    flow_box(s, 2.75, y, 9.75, 0.82, d, fill=WHITE, size=12.5)

# ==================== S10 原型展示 ====================
s = new_slide()
title_bar(s, "原型展示", "真实运行截图（Python + Streamlit 原型）")
picks = [("01_home", "首页双入口"), ("06_route_view", "路线画卷地图"),
         ("07_site_dialogue", "故地重游对话"), ("11_archive", "个人游历档案")]
shown = 0
for i, (name, label) in enumerate(picks):
    if not shot_available(name):
        continue
    x = 0.75 + (i % 2) * 6.2
    y = 1.65 + (i // 2) * 2.7
    s.shapes.add_picture(str(SHOT_DIR / f"{name}.png"), Inches(x), Inches(y),
                         width=Inches(5.9))
    add_box(s, x, y + 2.42, 5.9, 0.35, label, size=11, color=INK_L, align=PP_ALIGN.CENTER)
    shown += 1
if shown == 0:
    add_box(s, 2, 3, 9, 1, "（截图待录制：python scripts/record_demo.py）", size=18, color=INK_L)

# ==================== S11 创新点 ====================
s = new_slide()
title_bar(s, "项目创新点")
innovs = [
    ("① 一体化文旅业务闭环", "地点检索—偏好—成套路线—沉浸对话—打卡—合影—随笔—档案，一套链路深度打通"),
    ("② 双向融合路线生成", "用户个人喜好 × 历史人物真实生平爱好，定制专属成套行程，区别于普通旅游软件"),
    ("③ 第一人称叙事", "不只讲史实，更还原人物当时内心想法与情绪，让历史“活”起来"),
    ("④ 流式对话 + 语音输入", "AI 回复逐字浮现（st.write_stream）；浏览器原生语音识别（Web Speech，免 key）一句话自动发送"),
    ("⑤ 跨时代群聊", "多位不同朝代人物同席互辩，并发生成、完成一条浮现一条，隔离轮转防串台"),
    ("⑥ 双地图 + 高可用", "古风画卷 SVG + 实景交互地图（点击直达对话）；无 AI 时内置剧本降级，演示全链路永不中断"),
    ("⑦ 年轻化玩法", "古今人格测试、古人赠藏头诗、时空票根收藏——专业内容与年轻表达结合"),
]
for i, (t, d) in enumerate(innovs):
    y = 1.7 + i * 0.75
    flow_box(s, 0.9, y, 4.1, 0.6, t, fill=PAPER2, bold=True, size=14)
    add_box(s, 5.2, y + 0.05, 7.3, 0.6, d, size=13, color=INK_L)

# ==================== S12 不足与未来迭代 ====================
s = new_slide()
title_bar(s, "不足与未来迭代")
add_box(s, 0.9, 1.7, 5.5, 0.5, "现存不足", size=18, bold=True, color=SEAL)
add_box(s, 0.9, 2.3, 5.5, 3.2,
        [("· 人物立绘为程序生成的水墨占位图，待替换真实国风美术资源", {"size": 13}),
         ("· 合影为 PIL 静态合成，非生成式 AI 图像", {"size": 13}),
         ("· 语音输入依赖浏览器 Web Speech（Chrome/Edge 且需联网），离线 ASR 模型待接入", {"size": 13}),
         ("· 地点库 37 处故地，尚待扩充全国", {"size": 13}),
         ("· 实景交互地图为站点级点击直达，全国级动态地图未开发", {"size": 13})])
add_box(s, 6.9, 1.7, 5.5, 0.5, "未来迭代规划", size=18, bold=True, color=SEAL)
add_box(s, 6.9, 2.3, 5.5, 3.2,
        [("① 全国动态交互式地图：录入更多景点经纬度，实现跨城动态切换与轨迹动画（本次原型不开发）", {"size": 13}),
         ("② 动态数字人：接入 SadTalker 等方案，基于静态国风图片生成带口型表情的说话动画（本次原型不开发）", {"size": 13}),
         ("③ 离线 ASR 模型接入（funasr/faster-whisper）· ④ 扩充全国地点库与人物库 · ⑤ 移动端 PWA", {"size": 13})])

# ==================== S13 总结致谢 ====================
s = new_slide()
add_box(s, 1.5, 2.3, 10.3, 1.2,
        "千年晤旅——让每一次旅行，都有一位懂历史、有故事的故人同行。",
        size=28, bold=True, align=PP_ALIGN.CENTER)
add_box(s, 1.5, 3.6, 10.3, 0.6, "历史不应该是景点旁的冰冷说明牌，而是一场可以对话的重逢。",
        size=18, color=INK_L, align=PP_ALIGN.CENTER)
add_box(s, 1.5, 5.3, 10.3, 0.6, "感谢观看 · 恳请各位评委老师批评指正", size=20,
        color=SEAL, align=PP_ALIGN.CENTER)

# ==================== S14 AI辅助说明 ====================
s = new_slide()
title_bar(s, "AI辅助说明")
add_box(s, 0.9, 2.7, 11.5, 1.9,
        "本项目部分配图、文字表述由AI辅助制作，全部内容经过团队人工审核修改，"
        "核心功能代码、业务逻辑由团队自主完成。",
        size=17, align=PP_ALIGN.CENTER)
add_box(s, 0.9, 5.0, 11.5, 0.6, "千年晤旅项目组 · iCAN 大学生创新创业大赛", size=13,
        color=INK_L, align=PP_ALIGN.CENTER)

for i, sl in enumerate(prs.slides):
    if i > 0:
        footer(sl, i)

OUT = SUB / "千年晤旅_产品说明.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("PPT 已生成:", OUT, f"{OUT.stat().st_size/1024:.0f} KB · {len(prs.slides._sldIdLst)} 页")
